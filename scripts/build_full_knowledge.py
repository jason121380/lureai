#!/usr/bin/env python3
"""Extract the full Hair Brain source tree into private Markdown and RAG JSONL.

The source material is untrusted input. This script only reads documents and
never executes macros, formulas, links, or instructions found inside them.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import sqlite3
import subprocess
import tempfile
import zipfile
from collections import Counter
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable

PROJECT_ROOT = Path(__file__).resolve().parents[1]
SUPPORTED_OFFICE = {".ppt", ".pptx", ".xls", ".xlsx", ".doc", ".docx", ".pdf"}
IMAGE_TYPES = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".heic"}
PRIVATE_PATTERNS = [
    re.compile(r"https?://\S+", re.I),
    re.compile(r"[\w.+-]+@[\w.-]+\.[A-Za-z]{2,}"),
    re.compile(r"(?<!\d)(?:\+?886[- ]?)?0?9\d{2}[- ]?\d{3}[- ]?\d{3}(?!\d)"),
    re.compile(r"(?<!\d)(?:0\d{1,2}[- ]?)?\d{3,4}[- ]?\d{4}(?!\d)"),
    re.compile(r"(?<![A-Za-z0-9])[A-Z][12]\d{8}(?!\d)"),
    re.compile(r"(?<!\d)(?:\d[ -]?){15,19}(?!\d)"),
]
MENTION_PATTERN = re.compile(r"@[\w\-.\u3400-\u9fff]+")
ADDRESS_PATTERN = re.compile(
    r"(?<![\u3400-\u9fff])[\u3400-\u9fff]{2,3}[縣市]"
    r"[\u3400-\u9fff]{1,8}(?:區|鄉|鎮|市)[^\n，。；]{1,50}(?:號(?:之\d+)?(?:\d+樓)?|樓)"
)
DISTRICT_ADDRESS_PATTERN = re.compile(
    r"(?<![\u3400-\u9fff])[\u3400-\u9fff]{1,8}(?:區|鄉|鎮|市)"
    r"[\u3400-\u9fff]{1,12}(?:路|街|大道)[^\n，。；]{1,35}(?:號(?:之\d+)?(?:\d+樓)?|樓)"
)
HISTORICAL_NUMBER_PATTERN = re.compile(r"(?<![A-Za-z])\d[\d,]*(?:\.\d+)?%?")
SPACE_PATTERN = re.compile(r"[ \t]+")
EMPTY_LINES_PATTERN = re.compile(r"\n{3,}")
CUSTOMER_INCLUDE = re.compile(
    r"顧客|客戶|客服|客訴|服務流程|售後|接待|預約|髮型|頭型|臉型|毛髮|護髮|技術流程|現場開發|話術"
)
CUSTOMER_EXCLUDE = re.compile(
    r"業績|損益|員工|主管|管理ABC|組織|人才|活動企劃|廣告成效|目標設定|經營管理|店務|薪資|獎金|客資料|名單|人事|財務|離職|請假"
)
RAG_EXCLUDE = re.compile(
    r"客資料|名單|業績報表|設計師業績表|客數表|薪資|人事規章|財務|損益|離職|請假單|員工介紹|月行事曆"
)
DEPLOY_IDENTITY_FIELD = re.compile(
    r"((?:編號|姓名|店家|店名|店長|輔導人|電話|手機|地址|身分證(?:字號)?|帳號)\s*[:：]\s*)"
    r"([^|｜\n]{1,60})"
)


@dataclass
class Extracted:
    title: str
    sections: list[tuple[str, str]] = field(default_factory=list)
    status: str = "extracted"
    warnings: list[str] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)

    @property
    def text(self) -> str:
        return "\n\n".join(text for _, text in self.sections if text.strip())


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def clean_text(value: object) -> str:
    text = str(value or "").replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
    lines = [SPACE_PATTERN.sub(" ", line).strip() for line in text.splitlines()]
    return EMPTY_LINES_PATTERN.sub("\n\n", "\n".join(line for line in lines if line)).strip()


def shape_text(shape) -> list[str]:
    lines: list[str] = []
    if getattr(shape, "has_text_frame", False):
        value = clean_text(shape.text)
        if value:
            lines.append(value)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            values = [clean_text(cell.text).replace("\n", " / ") for cell in row.cells]
            if any(values):
                lines.append(" | ".join(values))
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            lines.extend(shape_text(child))
    return lines


def extract_pptx(path: Path) -> Extracted:
    from pptx import Presentation

    try:
        deck = Presentation(path)
    except (zipfile.BadZipFile, KeyError, ValueError, OSError) as exc:
        recovered = extract_pptx_xml(path)
        recovered.warnings.append(f"標準 PPTX 解析失敗，已從 XML 救回文字：{type(exc).__name__}: {exc}")
        return recovered
    sections = []
    for index, slide in enumerate(deck.slides, 1):
        values = []
        for shape in slide.shapes:
            values.extend(shape_text(shape))
        notes = ""
        try:
            notes = clean_text(slide.notes_slide.notes_text_frame.text)
        except Exception:
            pass
        if notes:
            values.append(f"講者備註：\n{notes}")
        sections.append((f"投影片 {index}", "\n\n".join(dict.fromkeys(v for v in values if v))))
    return Extracted(path.stem, sections, metadata={"slides": len(deck.slides)})


def extract_pptx_xml(path: Path) -> Extracted:
    import xml.etree.ElementTree as ET

    sections = []
    with zipfile.ZipFile(path) as archive:
        slide_names = sorted(
            (name for name in archive.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)),
            key=lambda name: int(re.search(r"(\d+)", Path(name).name).group(1)),
        )
        for index, name in enumerate(slide_names, 1):
            root = ET.fromstring(archive.read(name))
            values = [clean_text(node.text) for node in root.iter() if node.tag.endswith("}t") and clean_text(node.text)]
            sections.append((f"投影片 {index}", "\n".join(values)))
    return Extracted(path.stem, sections, metadata={"slides": len(sections), "recovered_from_xml": True})


def cell_value(cell) -> str:
    value = cell.value
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.isoformat(sep=" ")
    return clean_text(value).replace("\n", " / ")


def extract_xlsx(path: Path) -> Extracted:
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=False, read_only=True, keep_links=False)
    sections = []
    sheet_meta = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        nonempty = 0
        for row in sheet.iter_rows():
            pairs = []
            for cell in row:
                value = cell_value(cell)
                if value:
                    pairs.append(f"{cell.coordinate}={value}")
                    nonempty += 1
            if pairs:
                rows.append(" | ".join(pairs))
        sections.append((f"工作表：{sheet.title}", "\n".join(rows)))
        sheet_meta.append({"name": sheet.title, "nonempty_cells": nonempty})
    workbook.close()
    return Extracted(path.stem, sections, metadata={"sheets": sheet_meta})


def extract_docx(path: Path) -> Extracted:
    from docx import Document

    document = Document(path)
    sections: list[tuple[str, str]] = []
    current_title = "文件內容"
    current: list[str] = []
    for paragraph in document.paragraphs:
        text = clean_text(paragraph.text)
        if not text:
            continue
        style = str(getattr(paragraph.style, "name", ""))
        if style.lower().startswith("heading") and current:
            sections.append((current_title, "\n".join(current)))
            current_title = text
            current = []
        elif style.lower().startswith("heading"):
            current_title = text
        else:
            current.append(text)
    if current:
        sections.append((current_title, "\n".join(current)))
    for table_index, table in enumerate(document.tables, 1):
        rows = []
        for row in table.rows:
            values = [clean_text(cell.text).replace("\n", " / ") for cell in row.cells]
            if any(values):
                rows.append(" | ".join(values))
        sections.append((f"表格 {table_index}", "\n".join(rows)))
    return Extracted(path.stem, sections, metadata={"paragraphs": len(document.paragraphs), "tables": len(document.tables)})


def extract_pdf(path: Path) -> Extracted:
    from pypdf import PdfReader

    reader = PdfReader(path)
    if reader.is_encrypted:
        return Extracted(
            path.stem,
            [("受保護 PDF", "此 PDF 需要密碼，未嘗試繞過保護。")],
            status="protected",
            warnings=["PDF 已加密；需要檔案擁有者提供密碼才能抽取正文。"],
            metadata={"encrypted": True},
        )
    sections = []
    for index, page in enumerate(reader.pages, 1):
        text = clean_text(page.extract_text() or "")
        sections.append((f"頁面 {index}", text))
    result = Extracted(path.stem, sections, metadata={"pages": len(reader.pages)})
    if not result.text:
        result.status = "metadata_only"
        result.warnings.append("PDF 沒有可抽取文字，可能需要掃描件 OCR。")
    return result


def convert_legacy(path: Path, cache_root: Path, target_ext: str) -> Path:
    output_dir = cache_root / sha256(path)[:16]
    output_dir.mkdir(parents=True, exist_ok=True)
    expected = output_dir / f"{path.stem}{target_ext}"
    if expected.exists() and expected.stat().st_size:
        return expected
    command = ["soffice", "--headless", "--convert-to", target_ext.lstrip("."), "--outdir", str(output_dir), str(path)]
    completed = subprocess.run(command, capture_output=True, text=True, timeout=120)
    if completed.returncode != 0 or not expected.exists():
        detail = clean_text(completed.stderr or completed.stdout)
        raise RuntimeError(detail or "LibreOffice conversion failed")
    return expected


def extract_image(path: Path, ocr_binary: Path | None) -> Extracted:
    metadata = {}
    try:
        from PIL import Image
        with Image.open(path) as image:
            metadata = {"width": image.width, "height": image.height, "mode": image.mode}
    except Exception as exc:
        metadata = {"image_error": str(exc)}
    text = ""
    warnings = []
    if ocr_binary and ocr_binary.exists():
        completed = subprocess.run([str(ocr_binary), str(path)], capture_output=True, text=True, timeout=90)
        text = clean_text(completed.stdout)
        if completed.returncode != 0:
            warnings.append(clean_text(completed.stderr) or "OCR failed")
    else:
        warnings.append("找不到 OCR 工具。")
    status = "extracted" if text else "metadata_only"
    return Extracted(path.stem, [("圖片 OCR", text)], status=status, warnings=warnings, metadata=metadata)


def extract_database(path: Path) -> Extracted:
    sections = []
    metadata = {}
    try:
        connection = sqlite3.connect(f"file:{path}?mode=ro", uri=True)
        tables = [row[0] for row in connection.execute("SELECT name FROM sqlite_master WHERE type='table' ORDER BY name")]
        metadata["tables"] = tables
        for table in tables:
            safe_name = table.replace('"', '""')
            columns = [row[1] for row in connection.execute(f'PRAGMA table_info("{safe_name}")')]
            count = connection.execute(f'SELECT COUNT(*) FROM "{safe_name}"').fetchone()[0]
            sections.append((f"資料表：{table}", f"欄位：{', '.join(columns)}\n筆數：{count}"))
        connection.close()
        return Extracted(path.stem, sections, metadata=metadata)
    except Exception as exc:
        return Extracted(path.stem, [("資料庫資訊", "")], status="metadata_only", warnings=[str(exc)])


def extract_video(path: Path, video_ocr_binary: Path | None) -> Extracted:
    completed = subprocess.run(["mdls", "-name", "kMDItemDurationSeconds", "-name", "kMDItemPixelHeight", "-name", "kMDItemPixelWidth", str(path)], capture_output=True, text=True)
    metadata_text = clean_text(completed.stdout)
    frame_text = ""
    warnings = ["影片已做定時抽幀 OCR；音訊內容尚未轉成逐字稿。"]
    if video_ocr_binary and video_ocr_binary.exists():
        ocr = subprocess.run([str(video_ocr_binary), str(path)], capture_output=True, text=True, timeout=240)
        frame_text = clean_text(ocr.stdout)
        if ocr.returncode != 0:
            warnings.append(clean_text(ocr.stderr) or "影片抽幀 OCR 失敗。")
    return Extracted(
        path.stem,
        [("媒體資訊", metadata_text), ("定時抽幀 OCR", frame_text)],
        status="extracted" if frame_text else "metadata_only",
        warnings=warnings,
    )


def extract_json_summary(path: Path) -> Extracted:
    with path.open(encoding="utf-8") as handle:
        data = json.load(handle)
    if isinstance(data, dict) and isinstance(data.get("conversations"), list):
        text = (
            f"匯出時間：{data.get('exported_at', '')}\n"
            f"對話數：{data.get('conversation_count', len(data['conversations']))}\n"
            f"訊息數：{data.get('message_count', sum(len(c.get('messages', [])) for c in data['conversations']))}\n"
            "詳細內容已轉為去識別化對話案例 Markdown。"
        )
        return Extracted(path.stem, [("資料集摘要", text)], metadata={"dataset": "lurebot_conversations"})
    return Extracted(path.stem, [("JSON 摘要", f"頂層類型：{type(data).__name__}")], status="metadata_only")


def protected_or_corrupt(path: Path, error: Exception) -> Extracted:
    completed = subprocess.run(["file", "-b", str(path)], capture_output=True, text=True)
    signature = clean_text(completed.stdout)
    protected = "Security: 1" in signature
    label = "受保護檔案" if protected else "損壞或無法辨識的舊格式檔案"
    return Extracted(
        path.stem,
        [(label, "此檔案未取得可解析正文。")],
        status="protected" if protected else "corrupt_or_unknown",
        warnings=[f"{type(error).__name__}: {error}"],
        metadata={"file_signature": signature},
    )


def extract_file(path: Path, cache_root: Path, ocr_binary: Path | None, video_ocr_binary: Path | None) -> Extracted:
    suffix = path.suffix.lower()
    if path.name.startswith("~$"):
        return Extracted(path.stem, [("狀態", "Microsoft Office 暫存鎖定檔，沒有可用正文。")], status="skipped_temporary")
    if suffix == ".ppt":
        try:
            return extract_pptx(convert_legacy(path, cache_root, ".pptx"))
        except RuntimeError as exc:
            return protected_or_corrupt(path, exc)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".xls":
        try:
            return extract_xlsx(convert_legacy(path, cache_root, ".xlsx"))
        except RuntimeError as exc:
            return protected_or_corrupt(path, exc)
    if suffix == ".xlsx":
        return extract_xlsx(path)
    if suffix == ".doc":
        try:
            return extract_docx(convert_legacy(path, cache_root, ".docx"))
        except RuntimeError as exc:
            return protected_or_corrupt(path, exc)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix in IMAGE_TYPES:
        return extract_image(path, ocr_binary)
    if suffix == ".mp4":
        return extract_video(path, video_ocr_binary)
    if suffix == ".db":
        return extract_database(path)
    if suffix == ".json":
        return extract_json_summary(path)
    try:
        if path.read_bytes()[:3] == b"\xff\xd8\xff":
            return extract_image(path, ocr_binary)
    except OSError:
        pass
    return Extracted(path.stem, [("檔案資訊", "沒有適用的文字解析器。")], status="unsupported")


def summarize(extracted: Extracted) -> str:
    candidates = []
    for heading, text in extracted.sections:
        for value in [heading, *text.splitlines()]:
            value = clean_text(value).lstrip("-0123456789.、 ")
            if len(value) >= 6 and value not in candidates:
                candidates.append(value)
            if len(candidates) >= 5:
                break
        if len(candidates) >= 5:
            break
    if not candidates:
        return "此檔案目前只有中繼資料，尚未抽取到可索引正文。"
    return "；".join(candidates)[:600]


def markdown_for(relative: Path, source: Path, digest: str, extracted: Extracted, duplicate_of: str | None) -> str:
    frontmatter = {
        "source_path": relative.as_posix(),
        "source_type": source.suffix.lower().lstrip(".") or "unknown",
        "source_sha256": digest,
        "source_size": source.stat().st_size,
        "extract_status": extracted.status,
        "duplicate_of": duplicate_of or "",
    }
    lines = ["---"]
    lines.extend(f"{key}: {json.dumps(value, ensure_ascii=False)}" for key, value in frontmatter.items())
    lines.extend(["---", "", f"# {source.name}", "", "## 摘要", "", summarize(extracted), "", "## 中繼資料", ""])
    lines.append("```json")
    lines.append(json.dumps(extracted.metadata, ensure_ascii=False, indent=2))
    lines.append("```")
    if extracted.warnings:
        lines.extend(["", "## 抽取警告", ""])
        lines.extend(f"- {warning}" for warning in extracted.warnings)
    lines.extend(["", "## 詳細內容"])
    for heading, text in extracted.sections:
        lines.extend(["", f"### {heading}", "", text or "（沒有可抽取文字）"])
    return "\n".join(lines).rstrip() + "\n"


def safe_slug(value: str) -> str:
    return re.sub(r"[^A-Za-z0-9_-]+", "-", value).strip("-")[:80] or "conversation"


def collect_names(conversations: list[dict]) -> set[str]:
    names = set()
    for conversation in conversations:
        conversation_name = clean_text(conversation.get("conv_name", ""))
        if conversation_name:
            names.add(conversation_name)
            stripped = re.sub(r"^\[[^\]]+\]\s*", "", conversation_name)
            for segment in re.split(r"\s+-\s+", stripped):
                segment = segment.strip("-–—_()[] ")
                if len(segment) >= 2:
                    names.add(segment)
                for token in re.split(r"\s+|[/|]", segment):
                    token = token.strip("-–—_()[] ")
                    if len(token) >= 2:
                        names.add(token)
        for message in conversation.get("messages", []):
            sender = clean_text(message.get("sender", ""))
            if len(sender) >= 2:
                names.add(sender)
                for token in re.split(r"\s+|[-_/|]", sender):
                    token = token.strip()
                    if len(token) >= 2:
                        names.add(token)
    return names


def compile_name_redactor(names: set[str]) -> tuple[re.Pattern | None, list[str]]:
    ascii_names = sorted(
        (
            name for name in names
            if len(name) >= 3 and re.fullmatch(r"[A-Za-z0-9 ._/'-]+", name)
        ),
        key=len,
        reverse=True,
    )
    literal_names = sorted((name for name in names if name not in ascii_names), key=len, reverse=True)
    if not ascii_names:
        return None, literal_names
    alternatives = "|".join(re.escape(name) for name in ascii_names)
    return re.compile(rf"(?<![A-Za-z0-9])(?:{alternatives})(?![A-Za-z0-9])", re.I), literal_names


# Customer names often appear as surname + honorific ("吳小姐"); the sender
# name list never contains them, so mask the surname explicitly.
HONORIFIC_NAME_PATTERN = re.compile(r"([㐀-鿿])(小姐|先生|太太|女士)")
HONORIFIC_EXCLUDED_LEADS = set("這那各哪位小一二兩幾的是叫有大老她他我你妳您名")
# Wording that must never be treated as a personal name.
SANITIZE_KEEP = ("參考", "名留")


def _mask_honorific_names(value: str) -> str:
    def _replace(match: re.Match) -> str:
        if match.group(1) in HONORIFIC_EXCLUDED_LEADS:
            return match.group(0)
        return "[人名]" + match.group(2)

    return HONORIFIC_NAME_PATTERN.sub(_replace, value)


def sanitize_message(
    text: str,
    names: set[str],
    ascii_name_pattern: re.Pattern | None = None,
    literal_names: list[str] | None = None,
) -> str:
    value = clean_text(text)
    value = MENTION_PATTERN.sub("[提及對象]", value)
    value = ADDRESS_PATTERN.sub("[地址已移除]", value)
    value = DISTRICT_ADDRESS_PATTERN.sub("[地址已移除]", value)
    value = _mask_honorific_names(value)
    for keep in SANITIZE_KEEP:
        names.discard(keep)
    for pattern in PRIVATE_PATTERNS:
        value = pattern.sub("[敏感資訊已移除]", value)
    if literal_names is None:
        ascii_name_pattern, literal_names = compile_name_redactor(names)
    for name in literal_names:
        if len(name) >= 2:
            value = value.replace(name, "[人名]")
    if ascii_name_pattern:
        value = ascii_name_pattern.sub("[人名]", value)
    return clean_text(value)


def sanitize_deployable_text(text: str) -> str:
    value = sanitize_message(text, set())
    return DEPLOY_IDENTITY_FIELD.sub(lambda match: f"{match.group(1)}[已移除]", value)


def write_conversations(data: dict, output_root: Path) -> tuple[list[dict], list[dict]]:
    conversations = data.get("conversations", [])
    names = collect_names(conversations)
    ascii_name_pattern, literal_names = compile_name_redactor(names)
    markdown_root = output_root / "conversations"
    markdown_root.mkdir(parents=True, exist_ok=True)
    rag_rows = []
    manifest = []
    for index, conversation in enumerate(conversations, 1):
        original_id = str(conversation.get("conv_id", index))
        case_id = hashlib.sha256(original_id.encode()).hexdigest()[:16]
        messages = []
        for message in conversation.get("messages", []):
            text = sanitize_message(message.get("text", ""), names, ascii_name_pattern, literal_names)
            if not text or text in {"[提及對象]", "[人名]"}:
                continue
            messages.append({
                "role": "教練" if message.get("role") == "輔導" else "設計師",
                "month": str(message.get("sent_at", ""))[:7],
                "text": text,
            })
        md_path = markdown_root / f"case-{index:04d}-{case_id}.md"
        lines = [
            "---",
            f"case_id: {json.dumps(case_id)}",
            f"message_count: {len(messages)}",
            "access_level: \"internal_coaching\"",
            "deidentified: true",
            "---",
            "",
            f"# 去識別化輔導案例 {index:04d}",
            "",
            "## 摘要",
            "",
            f"此案例包含 {len(messages)} 則可用訊息，涵蓋設計師與教練的實際輔導往返。數字與方案均視為歷史案例，不代表現行標準。",
            "",
            "## 對話內容",
        ]
        for message_index, message in enumerate(messages, 1):
            lines.extend(["", f"### 訊息 {message_index}｜{message['role']}｜{message['month']}", "", message["text"]])
        md_path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")
        manifest.append({"case_id": case_id, "messages": len(messages), "markdown": str(md_path.relative_to(output_root))})

        window: list[dict] = []
        window_chars = 0
        chunk_index = 1
        for message in messages:
            rendered = f"{message['role']}（{message['month']}）：{message['text']}"
            if window and window_chars + len(rendered) > 1800:
                rag_rows.append(conversation_chunk(case_id, index, chunk_index, window))
                chunk_index += 1
                window = window[-2:]
                window_chars = sum(len(item["text"]) for item in window)
            window.append(message)
            window_chars += len(rendered)
        if window:
            rag_rows.append(conversation_chunk(case_id, index, chunk_index, window))
    return rag_rows, manifest


def conversation_chunk(case_id: str, case_index: int, chunk_index: int, messages: list[dict]) -> dict:
    rendered_messages = []
    for item in messages:
        historical_text = HISTORICAL_NUMBER_PATTERN.sub("[歷史數值]", item["text"])
        rendered_messages.append(f"{item['role']}（{item['month']}）：{historical_text}")
    text = (
        "【歷史輔導案例：僅供流程與判斷方法參考；數字、價格、時程、制度、活動與效果均非現行標準。】\n\n"
        + "\n\n".join(rendered_messages)
    )
    return {
        "chunk_id": f"coach-case:{case_id}:{chunk_index}",
        "doc_id": f"coach-case:{case_id}",
        "locator": f"case-{case_index:04d}-part-{chunk_index}",
        "section_title": f"去識別化輔導案例 {case_index:04d}",
        "text": text,
        "title": "設計師 1 對 1 歷史輔導案例",
        "source_file": f"private_sources/conversations/case-{case_index:04d}-{case_id}.md",
        "source_sha256": case_id,
        "category": "歷史輔導案例",
        "access_level": "internal_coaching",
        "rag_allowed": True,
        "review_status": "approved",
        "reviewer": "自動去識別化管線",
        "reviewed_at": datetime.now(timezone.utc).date().isoformat(),
        "historical_example": True,
    }


def split_text(text: str, max_chars: int = 1400) -> list[str]:
    paragraphs = [clean_text(value) for value in re.split(r"\n\s*\n|(?<=。)\s*", text) if clean_text(value)]
    chunks = []
    current = ""
    for paragraph in paragraphs:
        if len(paragraph) > max_chars:
            pieces = [paragraph[i:i + max_chars] for i in range(0, len(paragraph), max_chars)]
        else:
            pieces = [paragraph]
        for piece in pieces:
            if current and len(current) + len(piece) + 2 > max_chars:
                chunks.append(current)
                current = piece
            else:
                current = f"{current}\n\n{piece}".strip()
    if current:
        chunks.append(current)
    return chunks


def infer_category(relative: Path) -> str:
    value = relative.as_posix()
    rules = [
        (r"毛髮|護髮|技術|燙|染|髮型|頭型", "美髮技術"),
        (r"顧客|服務|售後|客訴|接待", "顧客服務"),
        (r"業績|損益|目標", "業績管理"),
        (r"主管|員工|人才|組織|職前", "人才與管理"),
        (r"網路|數位|廣告|社群|行銷", "數位行銷"),
        (r"店務|營運|店長|現場", "店務營運"),
    ]
    for pattern, category in rules:
        if re.search(pattern, value):
            return category
    return "企業知識"


def document_chunks(relative: Path, digest: str, extracted: Extracted, access_level: str) -> list[dict]:
    rows = []
    category = infer_category(relative)
    for section_index, (heading, text) in enumerate(extracted.sections, 1):
        for chunk_index, chunk in enumerate(split_text(text), 1):
            chunk = sanitize_deployable_text(chunk)
            if len(chunk) < 30:
                continue
            warning = "【歷史教材：僅供方法與流程參考；其中價格、時程、制度、效果與活動不得視為現行資訊。】"
            chunk = f"{warning}\n\n{chunk}"
            identity = hashlib.sha256(f"{relative}:{section_index}:{chunk_index}".encode()).hexdigest()[:20]
            rows.append({
                "chunk_id": f"source-doc:{identity}",
                "doc_id": f"source-doc:{digest[:16]}",
                "locator": f"section-{section_index}-part-{chunk_index}",
                "section_title": heading,
                "text": chunk,
                "title": relative.stem,
                "source_file": relative.as_posix(),
                "source_sha256": digest,
                "category": category,
                "access_level": access_level,
                "rag_allowed": True,
                "customer_service_allowed": access_level == "customer_service",
                "review_status": "approved",
                "reviewer": "自動抽取與安全分類管線",
                "reviewed_at": datetime.now(timezone.utc).date().isoformat(),
                "historical_source": True,
            })
    return rows


def write_jsonl(path: Path, rows: Iterable[dict]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as handle:
        for row in rows:
            handle.write(json.dumps(row, ensure_ascii=False) + "\n")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("source_root", type=Path)
    parser.add_argument("output_root", type=Path)
    parser.add_argument("--ocr-binary", type=Path)
    parser.add_argument("--video-ocr-binary", type=Path)
    parser.add_argument("--limit", type=int)
    args = parser.parse_args()

    source_root = args.source_root.resolve()
    output_root = args.output_root.resolve()
    extracted_root = output_root / "extracted"
    cache_root = output_root / "conversion_cache"
    extracted_root.mkdir(parents=True, exist_ok=True)
    cache_root.mkdir(parents=True, exist_ok=True)

    sources = sorted(path for path in source_root.rglob("*") if path.is_file() and path.name != ".DS_Store")
    if args.limit:
        sources = sources[:args.limit]
    digest_owner: dict[str, str] = {}
    manifest = []
    designer_rows = []
    customer_rows = []
    conversation_source: Path | None = None

    for index, source in enumerate(sources, 1):
        relative = source.relative_to(source_root)
        digest = sha256(source)
        duplicate_of = digest_owner.get(digest)
        digest_owner.setdefault(digest, relative.as_posix())
        md_path = extracted_root / relative.parent / f"{relative.name}.md"
        md_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            extracted = extract_file(source, cache_root, args.ocr_binary, args.video_ocr_binary)
        except Exception as exc:
            extracted = Extracted(source.stem, [("抽取錯誤", "")], status="failed", warnings=[f"{type(exc).__name__}: {exc}"])
        md_path.write_text(markdown_for(relative, source, digest, extracted, duplicate_of), encoding="utf-8")
        record = {
            "source": relative.as_posix(),
            "markdown": md_path.relative_to(output_root).as_posix(),
            "sha256": digest,
            "size": source.stat().st_size,
            "status": extracted.status,
            "duplicate_of": duplicate_of,
            "summary": summarize(extracted),
            "warnings": extracted.warnings,
            "sections": len(extracted.sections),
            "characters": len(extracted.text),
        }
        manifest.append(record)
        if source.name == "lurebot-conversations-20260831.json":
            conversation_source = source
        eligible_for_rag = not RAG_EXCLUDE.search(relative.as_posix())
        if not duplicate_of and extracted.status == "extracted" and source.suffix.lower() != ".json" and eligible_for_rag:
            designer_rows.extend(document_chunks(relative, digest, extracted, "internal_coaching"))
            if CUSTOMER_INCLUDE.search(relative.as_posix()) and not CUSTOMER_EXCLUDE.search(relative.as_posix()):
                customer_rows.extend(document_chunks(relative, digest, extracted, "customer_service"))
        print(json.dumps({"progress": f"{index}/{len(sources)}", "source": relative.as_posix(), "status": extracted.status}, ensure_ascii=False), flush=True)

    conversation_rows = []
    conversation_manifest = []
    if conversation_source:
        with conversation_source.open(encoding="utf-8") as handle:
            conversation_rows, conversation_manifest = write_conversations(json.load(handle), output_root)
        designer_rows.extend(conversation_rows)

    seed_coach = PROJECT_ROOT / "knowledge" / "designer_coaching_process.jsonl"
    seed_customer = PROJECT_ROOT / "knowledge" / "active_customer_service.jsonl"
    for seed, target in [(seed_coach, designer_rows), (seed_customer, customer_rows)]:
        if seed.exists():
            with seed.open(encoding="utf-8") as handle:
                target[:0] = [json.loads(line) for line in handle if line.strip()]

    write_jsonl(output_root / "rag" / "designer_coach_full.jsonl", designer_rows)
    write_jsonl(output_root / "rag" / "customer_service_full.jsonl", customer_rows)
    report = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "source_root": str(source_root),
        "source_files": len(sources),
        "unique_files": len(digest_owner),
        "status_counts": Counter(item["status"] for item in manifest),
        "markdown_files": len(manifest) + len(conversation_manifest),
        "conversation_cases": len(conversation_manifest),
        "conversation_chunks": len(conversation_rows),
        "designer_coach_chunks": len(designer_rows),
        "customer_service_chunks": len(customer_rows),
        "files": manifest,
        "conversations": conversation_manifest,
    }
    (output_root / "manifest.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({key: value for key, value in report.items() if key not in {"files", "conversations"}}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
